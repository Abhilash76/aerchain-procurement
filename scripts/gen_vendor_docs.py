import os
from docx import Document
from docx.shared import Inches
import pptx
from pptx.util import Inches as PptxInches
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import doc_utils

OUTPUT_DIR = "Vendor documents"

def gen_docx(filename, vendor_name, pages=35):
    print(f"Generating {filename} for {vendor_name}...")
    doc = Document()
    doc.add_heading(f'{vendor_name} - Technical Proposal', 0)
    
    for p in range(pages):
        page_num = p + 1
        doc.add_heading(f'Section {page_num}: {doc_utils.get_marketing_jargon()}', level=1)
        
        # Add page markers for AI to pick up (Source: Page X)
        doc.add_paragraph(f"--- INTERNAL MARKER: PAGE {page_num} ---")
        
        for _ in range(5):
            cost_cat = random.choice(list(doc_utils.VENDOR_PROFILES[vendor_name]["costs"].keys()))
            doc.add_paragraph(doc_utils.generate_messy_paragraph(vendor_name, include_cost=(p % 4 == 0), cost_category=cost_cat))
        
        if p % 5 == 0:
            doc.add_heading("Detailed Performance Metrics", level=2)
            table = doc.add_table(rows=1, cols=3)
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = 'Metric Area'
            hdr_cells[1].text = 'Rating (%)'
            hdr_cells[2].text = 'Benchmarking Insight'
            
            profile = doc_utils.VENDOR_PROFILES[vendor_name]
            for metric, score in profile["scores"].items():
                row_cells = table.add_row().cells
                row_cells[0].text = metric
                row_cells[1].text = str(score + random.randint(-5, 5))
                row_cells[2].text = doc_utils.get_marketing_jargon()
        
        doc.add_page_break()
    
    doc.save(os.path.join(OUTPUT_DIR, filename))

def gen_pptx(filename, vendor_name, slides=40):
    print(f"Generating {filename} for {vendor_name}...")
    prs = pptx.Presentation()
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"{vendor_name} Strategy Deck"
    subtitle.text = f"Confidential Proposal - {doc_utils.VENDOR_PROFILES[vendor_name]['style'].upper()}"
    
    for i in range(slides):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = f"Slide {i+1}: {doc_utils.get_marketing_jargon()}"
        tf = body_shape.text_frame
        tf.text = doc_utils.generate_messy_paragraph(vendor_name)
        
        # Add commercial info on some slides
        if i % 8 == 0:
            p = tf.add_paragraph()
            costs = doc_utils.VENDOR_PROFILES[vendor_name]["costs"]
            cat = random.choice(list(costs.keys()))
            p.text = f"PROPOSED INVESTMENT for {cat}: ${costs[cat]:,} (Ref: SOW-X-{i})"
            p.font.bold = True
        
    prs.save(os.path.join(OUTPUT_DIR, filename))

def gen_xlsx(filename, vendor_name, sheets=5):
    print(f"Generating {filename} for {vendor_name}...")
    wb = Workbook()
    profile = doc_utils.VENDOR_PROFILES[vendor_name]
    
    for s in range(sheets):
        if s == 0:
            ws = wb.active
            ws.title = "Commercial Summary"
        else:
            ws = wb.create_sheet(title=f"Detailed_Costs_{s}")
            
        ws.append(["Line Item", "Vendor Offering", "Unit Fee", "Quantity", "Subtotal", "Compliance Rating"])
        for cell in ws[1]:
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color="FFEE11", end_color="FFEE11", fill_type="solid")
            
        for r in range(2, 50):
            cat = random.choice(list(profile["costs"].keys()))
            score = profile["scores"].get(cat, 85)
            row = [
                cat,
                doc_utils.get_marketing_jargon(),
                profile["costs"][cat],
                1,
                "=C{0}*D{0}".format(r),
                f"{score}%"
            ]
            ws.append(row)
            
        ws.column_dimensions['B'].width = 40
        
    wb.save(os.path.join(OUTPUT_DIR, filename))

def gen_pdf(filename, vendor_name, pages=45):
    print(f"Generating {filename} for {vendor_name}...")
    c = canvas.Canvas(os.path.join(OUTPUT_DIR, filename), pagesize=letter)
    width, height = letter
    profile = doc_utils.VENDOR_PROFILES[vendor_name]
    
    for p in range(pages):
        c.setFont("Helvetica-Bold", 14)
        c.drawString(72, height - 72, f"{vendor_name} - [OFFICIAL PROPOSAL PAGE {p+1}]")
        c.setFont("Helvetica", 9)
        
        y = height - 100
        for _ in range(15):
            include_cost = (p % 6 == 0)
            cost_cat = random.choice(list(profile["costs"].keys()))
            text = doc_utils.generate_messy_paragraph(vendor_name, include_cost=include_cost, cost_category=cost_cat)
            
            words = text.split()
            lines = [" ".join(words[i:i+12]) for i in range(0, len(words), 12)]
            for line in lines:
                c.drawString(72, y, line)
                y -= 11
                if y < 72: break
            y -= 15
            if y < 72: break
            
        c.showPage()
    c.save()

import random # Ensure random is available if not already in doc_utils scope

if __name__ == "__main__":
    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)
        
    # Generate distinct proposals for each vendor
    gen_docx("Nexus_Technical_Proposal.docx", "Nexus Creative", pages=35)
    gen_pptx("Global_Marketing_Deck.pptx", "Global Media Hub", slides=40)
    gen_xlsx("Velocity_Commercial_Bid.xlsx", "Velocity Studios", sheets=3)
    gen_pdf("Nexus_Compliance_Addendum.pdf", "Nexus Creative", pages=20)
    gen_pdf("Master_Vendor_DeepDive.pdf", "Global Media Hub", pages=45)
    
    print("All multi-vendor test documents generated successfully.")
